"""Generate PedaForge project proposal slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(0x01, 0xAB, 0xAA)
PRIMARY_DARK = RGBColor(0x00, 0x8C, 0x8B)
SECONDARY = RGBColor(0x2D, 0x2A, 0x5E)
ACCENT = RGBColor(0xFE, 0xCD, 0x1B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xFA, 0xF6, 0xF0)
TEXT_DARK = RGBColor(0x3D, 0x3E, 0x3F)
TEXT_LIGHT = RGBColor(0x6B, 0x72, 0x80)
DANGER = RGBColor(0xE8, 0x06, 0x3C)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)
WARNING = RGBColor(0xFF, 0x9E, 0x18)
INFO = RGBColor(0x3B, 0x82, 0xF6)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, size=18,
                 bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT,
                 font_name="Arial"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_bullet_list(slide, left, top, width, height, items, size=14,
                    color=TEXT_DARK):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return txbox


def add_card(slide, left, top, width, height, title, body, accent_color,
             title_size=14, body_size=11):
    card = add_shape(slide, left, top, width, height, WHITE)

    add_text_box(slide, left + Inches(0.2), top + Inches(0.1),
                 width - Inches(0.3), Inches(0.4),
                 title, size=title_size, bold=True, color=SECONDARY)

    add_text_box(slide, left + Inches(0.2), top + Inches(0.45),
                 width - Inches(0.3), height - Inches(0.55),
                 body, size=body_size, color=TEXT_LIGHT)


def slide_title_bar(slide):
    bar = add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), PRIMARY)
    return bar


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ─── SLIDE 1: Title ─────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)

    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625),
              RGBColor(0x1A, 0x15, 0x35))

    add_shape(slide, Inches(0), Inches(4.8), Inches(10), Inches(0.06), PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.5),
                 "ECDA Early Childhood Innovation Sandbox Proposal", size=14,
                 color=RGBColor(0xB8, 0xA0, 0xD0), font_name="Arial")

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(1.2),
                 "PedaForge", size=44, bold=True, color=WHITE)

    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.8),
                 "AI-Powered Pedagogical Development Platform\n"
                 "for Singapore Early Childhood Education",
                 size=20, color=RGBColor(0xC8, 0xB8, 0xE0))

    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(8.4), Inches(0.4),
                 "Aligned with ECDA Industry Digital Plan (IDP 2.0)",
                 size=13, color=RGBColor(0xB8, 0xA0, 0xD0))

    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(4), Inches(0.4),
                 "Duration: 6 Months  |  Budget: SGD $20,000",
                 size=12, color=RGBColor(0x8B, 0x80, 0xA8))

    add_text_box(slide, Inches(5.5), Inches(4.2), Inches(4), Inches(0.4),
                 "May 2026", size=12,
                 color=RGBColor(0x8B, 0x80, 0xA8), alignment=PP_ALIGN.RIGHT)

    # ─── SLIDE 2: The Problem ───────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.35),
                 "THE CHALLENGE", size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.5),
                 "Singapore's ECE Documentation Crisis",
                 size=26, bold=True, color=SECONDARY)

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.5),
                 "Excessive administrative burden is directly harming pedagogical "
                 "quality and educator retention in the early childhood sector.",
                 size=13, color=TEXT_LIGHT)

    cards = [
        ("Emotional Exhaustion", "8 hrs/week", "on documentation",
         "Administrative compliance tasks cause severe emotional exhaustion. "
         "Educators spend evenings and weekends on paperwork.", DANGER),
        ("Pedagogical Compromise", "35% less", "interaction time",
         "Cognitive overload from multi-framework documentation forces "
         "one-size-fits-all planning, reducing individualised learning.", WARNING),
        ("Siloed Leadership", "Zero", "real-time visibility",
         "Directors lack data on daily pedagogical quality. Annual "
         "appraisals are subjective and disconnected from the Skills Framework.", INFO),
    ]

    for i, (title, stat, stat_label, desc, color) in enumerate(cards):
        x = Inches(0.6 + i * 3.1)
        card_bg = add_shape(slide, x, Inches(1.8), Inches(2.8), Inches(3.4),
                            LIGHT_BG)

        add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(2.4),
                     Inches(0.3), title, size=13, bold=True, color=SECONDARY)

        add_text_box(slide, x + Inches(0.2), Inches(2.4), Inches(2.4),
                     Inches(0.6), stat, size=32, bold=True, color=color)

        add_text_box(slide, x + Inches(0.2), Inches(2.95), Inches(2.4),
                     Inches(0.3), stat_label, size=11, color=TEXT_LIGHT)

        add_text_box(slide, x + Inches(0.2), Inches(3.4), Inches(2.4),
                     Inches(1.5), desc, size=10, color=TEXT_LIGHT)

    add_text_box(slide, Inches(0.6), Inches(5.25), Inches(8.8), Inches(0.3),
                 "Source: OECD TALIS Starting Strong (2024), "
                 "Early Childhood Research Quarterly (2024)",
                 size=8, color=TEXT_LIGHT, alignment=PP_ALIGN.RIGHT)

    # ─── SLIDE 3: Solution Overview ─────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.35),
                 "THE SOLUTION", size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.5),
                 "Four Integrated Modules, One Ecosystem",
                 size=26, bold=True, color=SECONDARY)

    modules = [
        ("1. Smart Lesson Planner",
         "AI generates differentiated scaffolding based on "
         "classroom learning profiles. Maps to EYDF, NEL, "
         "and implements Rehearse & Retrieve scheduling.",
         "EYDF | NEL | Differentiation", PRIMARY),
        ("2. Portfolio & Profiling",
         "Capture observations and work samples. AI drafts "
         "developmental narratives. 'Living Profile' feeds "
         "back into lesson planning.",
         "Living Profile | Assessment", INFO),
        ("3. AI Coaching Agent",
         "Four coaching modes: Reflective Practice, QTT "
         "Deep Dive, Socratic Inquiry, Scenario Analysis. "
         "Anti-sycophancy directives built in.",
         "QTT | Skills Framework", ACCENT),
        ("4. Director Dashboard",
         "Centre-wide QTT averages, educator profiles, "
         "PLC collaboration topics, and SFw-aligned "
         "professional development recommendations.",
         "LNA | SFw | IDP", RGBColor(0x8B, 0x5C, 0xF6)),
    ]

    for i, (title, desc, tags, color) in enumerate(modules):
        row = i // 2
        col = i % 2
        x = Inches(0.6 + col * 4.6)
        y = Inches(1.2 + row * 2.1)

        card_bg = add_shape(slide, x, y, Inches(4.2), Inches(1.9), LIGHT_BG)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.8),
                     Inches(0.35), title, size=14, bold=True, color=SECONDARY)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.8),
                     Inches(0.9), desc, size=10, color=TEXT_LIGHT)

        add_text_box(slide, x + Inches(0.2), y + Inches(1.45), Inches(3.8),
                     Inches(0.3), tags, size=8, bold=True, color=color)

    # ─── SLIDE 4: AI Coaching Deep Dive ─────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)

    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), ACCENT)

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.35),
                 "SIGNATURE FEATURE", size=10, bold=True, color=ACCENT)

    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.5),
                 "AI Coaching Engine: Four Modes, One Growth Journey",
                 size=24, bold=True, color=WHITE)

    modes = [
        ("Reflective Practice",
         "Gibbs' Reflective Cycle\n\nGentle, structured reflection "
         "on daily teaching episodes.\n\nBest for: Daily reflection",
         "Supportive", 1),
        ("QTT Deep Dive",
         "ECDA Quality Teaching Tool\n\nTargeted coaching aligned "
         "to specific QTT rubric domains.\n\nBest for: Improvement",
         "Focused", 2),
        ("Socratic Inquiry",
         "Six question types\n\nProbing questions that surface "
         "hidden pedagogical assumptions.\n\nBest for: Assumptions",
         "Challenging", 3),
        ("Scenario Analysis",
         "Gary Klein Pre-Mortem\n\nStress-test lesson plans "
         "before classroom execution.\n\nBest for: Preparation",
         "Rigorous", 4),
    ]

    for i, (title, desc, intensity, bars) in enumerate(modes):
        x = Inches(0.4 + i * 2.35)
        y = Inches(1.3)

        card_bg = add_shape(slide, x, y, Inches(2.15), Inches(2.8),
                            RGBColor(0x3D, 0x26, 0x63))

        add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(1.85),
                     Inches(0.35), title, size=12, bold=True, color=WHITE)

        add_text_box(slide, x + Inches(0.15), y + Inches(0.5), Inches(1.85),
                     Inches(1.6), desc, size=9, color=RGBColor(0xB0, 0xA0, 0xC8))

        bar_y = y + Inches(2.2)
        for b in range(4):
            bar_color = ACCENT if b < bars else RGBColor(0x3A, 0x30, 0x60)
            add_shape(slide, x + Inches(0.15 + b * 0.45), bar_y,
                      Inches(0.35), Inches(0.08), bar_color)

        add_text_box(slide, x + Inches(0.15), bar_y + Inches(0.15),
                     Inches(1.85), Inches(0.25), intensity,
                     size=8, color=RGBColor(0xB8, 0xA0, 0xD0))

    safeguards = [
        "Anti-sycophancy: No reflexive validation, specific earned praise only",
        "Safety guardrails: Auto de-escalation on distress signals",
        "QTT tagging: Every coaching response tagged to QTT domains",
        "Referral pathway: 'Speak to a Real Coach' at any time",
    ]
    for i, text in enumerate(safeguards):
        add_text_box(slide, Inches(0.6), Inches(4.3 + i * 0.28),
                     Inches(8.8), Inches(0.25),
                     f"  {text}", size=9,
                     color=RGBColor(0xB8, 0xA0, 0xD0))

    # ─── SLIDE 5: Ecosystem Flow ────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.35),
                 "HOW IT WORKS", size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.5),
                 "The PedaForge Synergistic Ecosystem",
                 size=26, bold=True, color=SECONDARY)

    steps = [
        ("1. Observe", "Educator captures\nobservations, work\nsamples, anecdotes",
         PRIMARY),
        ("2. Profile", "AI builds Living\nProfile: strengths,\nneeds, ZPD",
         INFO),
        ("3. Plan", "AI generates\ndifferentiated\nlesson plans",
         RGBColor(0x8B, 0x5C, 0xF6)),
        ("4. Teach", "Educator delivers\nwith scaffolding\nstrategies",
         ACCENT),
        ("5. Reflect", "AI Coach guides\npedagogical\nreflection",
         DANGER),
        ("6. Grow", "Director sees\ninsights, assigns\ntargeted PD",
         SUCCESS),
    ]

    for i, (title, desc, color) in enumerate(steps):
        x = Inches(0.3 + i * 1.57)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.3), Inches(1.3),
            Inches(0.9), Inches(0.9))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        add_text_box(slide, x + Inches(0.3), Inches(1.45), Inches(0.9),
                     Inches(0.7), title.split(". ")[0] + ".",
                     size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x, Inches(2.35), Inches(1.5), Inches(0.3),
                     title.split(". ")[1], size=12, bold=True,
                     color=SECONDARY, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x, Inches(2.65), Inches(1.5), Inches(1.0),
                     desc, size=9, color=TEXT_LIGHT,
                     alignment=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            add_text_box(slide, x + Inches(1.3), Inches(1.55),
                         Inches(0.4), Inches(0.5), "→",
                         size=24, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    feedback_bg = add_shape(slide, Inches(0.6), Inches(3.9), Inches(8.8),
                            Inches(0.7), LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(3.95), Inches(8.4), Inches(0.6),
                 "Continuous Feedback Loop: Portfolio observations drive lesson "
                 "planning → Lesson execution feeds coaching reflections → "
                 "Coaching insights inform director's PD decisions → "
                 "Growth targets cycle back to observations",
                 size=10, color=TEXT_LIGHT)

    add_text_box(slide, Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.5),
                 "Children's Learning Outcomes + Educator Professional Growth "
                 "= Quality ECE", size=13, bold=True, color=PRIMARY,
                 alignment=PP_ALIGN.CENTER)

    # ─── SLIDE 6: Architecture ──────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.35),
                 "TECHNICAL DESIGN", size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.5),
                 "Enterprise-Grade Architecture",
                 size=26, bold=True, color=SECONDARY)

    stack = [
        ("Frontend", "React / Next.js, TypeScript, TailwindCSS", INFO),
        ("Backend", "Python 3.11+, FastAPI, Uvicorn", PRIMARY),
        ("AI / LLM", "Claude Haiku 4.5 via Azure AI Foundry", ACCENT),
        ("Database", "PostgreSQL (children, plans, portfolios)", RGBColor(0x8B, 0x5C, 0xF6)),
        ("Auth", "JWT + bcrypt, RBAC (educator / director / admin)", DANGER),
        ("Hosting", "Azure App Service (SG region), Key Vault", SECONDARY),
    ]

    for i, (layer, tech, color) in enumerate(stack):
        y = Inches(1.15 + i * 0.58)
        add_shape(slide, Inches(0.6), y, Inches(4.2), Inches(0.48), LIGHT_BG)

        add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(1.3),
                     Inches(0.35), layer, size=11, bold=True, color=SECONDARY)

        add_text_box(slide, Inches(2.2), y + Inches(0.05), Inches(2.5),
                     Inches(0.35), tech, size=10, color=TEXT_LIGHT)

    compliance_items = [
        ("PDPA Compliant", "Zero public AI training; child data never used "
         "for model training"),
        ("SG Data Residency", "All data hosted in Singapore Azure region "
         "(southeastasia)"),
        ("Role-Based Access", "Educators see only assigned cohorts; Directors "
         "get aggregated centre view"),
        ("Encryption", "AES-256 at rest, TLS 1.2+ in transit, "
         "Key Vault for secrets"),
        ("Audit Trail", "All actions logged with user ID and timestamp"),
    ]

    add_text_box(slide, Inches(5.4), Inches(1.15), Inches(4.2), Inches(0.35),
                 "Data Privacy & Compliance", size=14, bold=True,
                 color=SECONDARY)

    for i, (title, desc) in enumerate(compliance_items):
        y = Inches(1.6 + i * 0.7)
        add_text_box(slide, Inches(5.4), y, Inches(4.2), Inches(0.25),
                     title, size=10, bold=True, color=PRIMARY)
        add_text_box(slide, Inches(5.4), y + Inches(0.22), Inches(4.2),
                     Inches(0.4), desc, size=9, color=TEXT_LIGHT)

    # ─── SLIDE 7: KPIs ─────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.35),
                 "IMPACT TARGETS", size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.5),
                 "Measurable Outcomes", size=26, bold=True, color=SECONDARY)

    kpis = [
        ("60%", "Documentation\nTime Saved", "8 hrs to 3 hrs/week", PRIMARY),
        ("100%", "Differentiated\nLesson Plans", "2+ profiles per plan", INFO),
        ("85%", "IDP Accuracy\nRating", "Educator agreement", SUCCESS),
        ("2x", "Weekly Coaching\nCycles", "Per educator", ACCENT),
        ("+0.5", "QTT Score\nImprovement", "Centre average", RGBColor(0x8B, 0x5C, 0xF6)),
        ("90%", "Educator\nRetention Intent", "Would continue", DANGER),
    ]

    for i, (value, label, sub, color) in enumerate(kpis):
        row = i // 3
        col = i % 3
        x = Inches(0.6 + col * 3.1)
        y = Inches(1.2 + row * 2.0)

        card_bg = add_shape(slide, x, y, Inches(2.8), Inches(1.8), LIGHT_BG)

        add_text_box(slide, x, y + Inches(0.15), Inches(2.8), Inches(0.7),
                     value, size=36, bold=True, color=color,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x, y + Inches(0.85), Inches(2.8), Inches(0.5),
                     label, size=12, bold=True, color=SECONDARY,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x, y + Inches(1.35), Inches(2.8), Inches(0.3),
                     sub, size=9, color=TEXT_LIGHT,
                     alignment=PP_ALIGN.CENTER)

    # ─── SLIDE 8: Timeline ──────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.35),
                 "6-MONTH ROADMAP", size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.5),
                 "Sandbox Implementation Timeline",
                 size=26, bold=True, color=SECONDARY)

    phases = [
        ("Months 1–2", "Development &\nAI Foundation",
         "Infrastructure, AI prompts,\nplatform build, PDPA", PRIMARY),
        ("Month 3", "Pilot\nDeployment",
         "2–3 centres onboarded,\neducator workshops", INFO),
        ("Months 4–5", "Active\nExecution",
         "Live classroom usage,\nfirst IDPs generated", ACCENT),
        ("Month 6", "Impact Assessment\n& Reporting",
         "Findings, commercialisation\nroadmap, ECDA presentation", DANGER),
    ]

    timeline_y = Inches(1.4)
    add_shape(slide, Inches(0.6), timeline_y + Inches(0.45),
              Inches(8.8), Inches(0.04), RGBColor(0xE2, 0xE8, 0xF0))

    for i, (period, title, desc, color) in enumerate(phases):
        x = Inches(0.4 + i * 1.55)

        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.6), timeline_y + Inches(0.3),
            Inches(0.32), Inches(0.32))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        add_text_box(slide, x, timeline_y - Inches(0.15), Inches(1.5),
                     Inches(0.3), period, size=8, bold=True, color=color,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x, timeline_y + Inches(0.8), Inches(1.5),
                     Inches(0.5), title, size=10, bold=True,
                     color=SECONDARY, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x, timeline_y + Inches(1.4), Inches(1.5),
                     Inches(0.7), desc, size=8, color=TEXT_LIGHT,
                     alignment=PP_ALIGN.CENTER)

    milestones = [
        "Month 2: AI prompts validated against QTT rubric",
        "Month 3: First educators onboarded with training",
        "Month 5: First SFw-aligned IDPs auto-generated",
        "Month 6: Final impact report to ECDA",
    ]

    add_text_box(slide, Inches(0.6), Inches(3.9), Inches(4), Inches(0.3),
                 "Key Milestones", size=12, bold=True, color=SECONDARY)

    for i, ms in enumerate(milestones):
        add_text_box(slide, Inches(0.8), Inches(4.25 + i * 0.28),
                     Inches(8.6), Inches(0.25),
                     f"✓  {ms}", size=9, color=TEXT_LIGHT)

    # ─── SLIDE 9: Budget ────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.35),
                 "INVESTMENT", size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.5),
                 "Budget Breakdown (SGD)", size=26, bold=True, color=SECONDARY)

    budget_phases = [
        ("Development & Infrastructure", "$9,800", "49%", PRIMARY),
        ("Security & Compliance", "$3,200", "16%", INFO),
        ("Change Management", "$4,400", "22%", ACCENT),
        ("Evaluation & Reporting", "$2,600", "13%",
         RGBColor(0x8B, 0x5C, 0xF6)),
    ]

    for i, (phase, amount, pct, color) in enumerate(budget_phases):
        y = Inches(1.3 + i * 0.75)

        bar_width = float(pct.replace("%", "")) / 100 * 5.5

        add_shape(slide, Inches(0.6), y, Inches(5.5), Inches(0.55),
                  LIGHT_BG)
        add_shape(slide, Inches(0.6), y, Inches(bar_width), Inches(0.55),
                  color)

        add_text_box(slide, Inches(0.75), y + Inches(0.08), Inches(3),
                     Inches(0.35), phase, size=10, bold=True, color=WHITE)

        add_text_box(slide, Inches(6.3), y + Inches(0.08), Inches(1.2),
                     Inches(0.35), amount, size=12, bold=True,
                     color=SECONDARY, alignment=PP_ALIGN.RIGHT)

        add_text_box(slide, Inches(7.6), y + Inches(0.08), Inches(0.8),
                     Inches(0.35), pct, size=10, color=TEXT_LIGHT)

    total_bg = add_shape(slide, Inches(0.6), Inches(4.4), Inches(8.2),
                         Inches(0.6), SECONDARY)
    add_text_box(slide, Inches(0.8), Inches(4.45), Inches(4), Inches(0.45),
                 "Total Sandbox Funding Required", size=14, bold=True,
                 color=WHITE)
    add_text_box(slide, Inches(5.5), Inches(4.45), Inches(3.1), Inches(0.45),
                 "SGD $20,000", size=20, bold=True, color=ACCENT,
                 alignment=PP_ALIGN.RIGHT)

    details_left = [
        "Full-Stack + AI Dev: $500/day x 16 days",
        "Azure Hosting: $200/month x 6 months",
        "Claude Haiku via Azure AI Foundry: $100/month x 6 months",
    ]

    details_right = [
        "PDPA + Security: $3,200",
        "Pedagogical Consultant: $600/day x 4 days",
        "Educator Workshops: $2,000",
        "Evaluation + Reporting: $2,600",
    ]

    for col, items in enumerate([details_left, details_right]):
        x = Inches(0.6 + col * 4.5)
        for i, item in enumerate(items):
            add_text_box(slide, x, Inches(5.15) - Inches(0.05) + Inches(i * 0.22),
                         Inches(4.2), Inches(0.2),
                         item, size=8, color=TEXT_LIGHT)

    # ─── SLIDE 10: Team ─────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.35),
                 "THE TEAM", size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.5),
                 "Cross-Disciplinary Expertise",
                 size=26, bold=True, color=SECONDARY)

    team = [
        ("Project Lead", "0.3 FTE",
         "Timeline, budget, stakeholder liaison,\nECDA comms, evaluation & reporting",
         PRIMARY),
        ("Technical Lead", "1.0 FTE",
         "Full-stack development, AI integration\nvia Azure AI Foundry, deployment",
         INFO),
        ("Pedagogical Specialist", "0.2 FTE",
         "EYDF, NEL, QTT, SFw validation,\neducator training facilitation",
         ACCENT),
    ]

    for i, (role, fte, desc, color) in enumerate(team):
        x = Inches(0.6 + i * 3.1)
        y = Inches(1.2)

        card_bg = add_shape(slide, x, y, Inches(2.8), Inches(1.8), LIGHT_BG)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1.05), y + Inches(0.15),
            Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        add_text_box(slide, x, y + Inches(0.9), Inches(2.8), Inches(0.3),
                     role, size=12, bold=True, color=SECONDARY,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x, y + Inches(1.15), Inches(2.8), Inches(0.2),
                     fte, size=9, color=color, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + Inches(0.15), y + Inches(1.35), Inches(2.5),
                     Inches(0.4), desc, size=8, color=TEXT_LIGHT,
                     alignment=PP_ALIGN.CENTER)

    # ─── SLIDE 11: Framework Alignment ──────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.35),
                 "ALIGNMENT", size=10, bold=True, color=PRIMARY)

    add_text_box(slide, Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.5),
                 "Grounded in Singapore's ECE Frameworks",
                 size=26, bold=True, color=SECONDARY)

    frameworks = [
        ("ECDA IDP 2.0",
         "Direct alignment with Early Childhood Innovation Sandbox objectives. "
         "AI-enabled solution for sector-wide transformation.",
         "Strategic Fit", PRIMARY),
        ("EYDF",
         "Early Years Development Framework. All learning outcomes "
         "mapped and tagged in lesson plans and portfolios.",
         "Curriculum", INFO),
        ("NEL",
         "Nurturing Early Learners. Activity design and assessment "
         "indicators aligned to NEL learning areas.",
         "Learning", SUCCESS),
        ("ECDA QTT",
         "Quality Teaching Tool. Coaching engine built around QTT "
         "rubric domains. Director dashboard tracks QTT scores.",
         "Quality", ACCENT),
        ("SFw for ECCE",
         "Skills Framework for Early Childhood. Professional "
         "development mapped to TSC codes. WSQ course recommendations.",
         "Growth", RGBColor(0x8B, 0x5C, 0xF6)),
        ("PDPA",
         "Personal Data Protection Act. Zero public AI training. "
         "SG data residency. RBAC. Encrypted storage.",
         "Compliance", DANGER),
    ]

    for i, (name, desc, category, color) in enumerate(frameworks):
        row = i // 3
        col = i % 3
        x = Inches(0.6 + col * 3.1)
        y = Inches(1.2 + row * 2.0)

        card_bg = add_shape(slide, x, y, Inches(2.8), Inches(1.8), LIGHT_BG)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(2.4),
                     Inches(0.3), name, size=14, bold=True, color=SECONDARY)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.45), Inches(2.4),
                     Inches(0.3), category, size=9, bold=True, color=color)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.75), Inches(2.4),
                     Inches(0.9), desc, size=9, color=TEXT_LIGHT)

    # ─── SLIDE 12: Call to Action ───────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECONDARY)

    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(5.625),
              RGBColor(0x1A, 0x15, 0x35))
    add_shape(slide, Inches(0), Inches(4.8), Inches(10), Inches(0.06), PRIMARY)

    add_text_box(slide, Inches(1), Inches(1.0), Inches(8), Inches(0.8),
                 "PedaForge", size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(1.8), Inches(8), Inches(0.6),
                 "Give educators back the time that matters most.",
                 size=20, color=RGBColor(0xC8, 0xB8, 0xE0),
                 alignment=PP_ALIGN.CENTER)

    highlights = [
        "60% reduction in documentation time",
        "AI coaching aligned to ECDA Quality Teaching Tool",
        "Profile-driven differentiated lesson planning",
        "Real-time director insights for targeted professional development",
    ]

    for i, h in enumerate(highlights):
        add_text_box(slide, Inches(2), Inches(2.7 + i * 0.35), Inches(6),
                     Inches(0.3), f"✓  {h}", size=13,
                     color=RGBColor(0xB8, 0xA0, 0xD0),
                     alignment=PP_ALIGN.CENTER)

    cta_bg = add_shape(slide, Inches(3), Inches(4.2), Inches(4),
                       Inches(0.5), PRIMARY)
    add_text_box(slide, Inches(3), Inches(4.22), Inches(4), Inches(0.45),
                 "SGD $20,000  |  6 Months  |  IDP 2.0 Aligned",
                 size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    output_path = "/home/dmgadmin/sandbox/pedaforge/slides/PedaForge_Proposal.pptx"
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = create_presentation()
    print(f"Presentation saved to: {path}")
